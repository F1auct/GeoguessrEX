"""
Add 4 "需求分析" slides to front of v2.0.pptx, matching existing style.
"""
import sys; sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

INK   = RGBColor(0x44,0x37,0x28)
CARD  = RGBColor(0xEB,0xE2,0xE0)
ACC   = RGBColor(0xB4,0x4D,0x28)
GRN   = RGBColor(0x24,0x4C,0x47)
WHITE = RGBColor(0xFF,0xF8,0xEF)
MUTED = RGBColor(0x83,0x5E,0x54)
FONT_H = '黑体'
FONT_B = 'Open Sans'
SLIDE_W, SLIDE_H = 9144000, 5143500
LM = Emu(400000)
CW = SLIDE_W - LM*2

prs = Presentation('v2.0.pptx')
layout = prs.slide_layouts[0]

def rect(s, l, t, w, h, fill=None, radius=0.08):
    sh = s.shapes.add_shape(5, l, t, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else: sh.fill.background()
    sh.line.fill.background(); sh.adjustments[0] = radius
    return sh

def tb(s, l, t, w, h, text, font=FONT_B, size=Pt(20), color=INK, bold=False, align=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h); bx.word_wrap = True
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.name = font; p.font.size = size; p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return bx, tf

def card(s, l, t, w, h, title, lines, title_color=ACC, title_size=Pt(18)):
    sh = rect(s, l, t, w, h, fill=CARD)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left=Emu(140000); tf.margin_right=Emu(140000); tf.margin_top=Emu(90000); tf.margin_bottom=Emu(90000)
    p = tf.paragraphs[0]; p.text = title
    p.font.name = FONT_H; p.font.size = title_size; p.font.color.rgb = title_color; p.font.bold = True; p.space_after = Pt(10)
    for ln in lines:
        p = tf.add_paragraph(); p.text = ln
        p.font.name = FONT_B; p.font.size = Pt(14); p.font.color.rgb = INK; p.space_after = Pt(6)
    return sh

def colored_card(s, l, t, w, h, title, desc, color):
    sh = rect(s, l, t, w, h, fill=color)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left=Emu(80000); tf.margin_right=Emu(80000); tf.margin_top=Emu(50000)
    p = tf.paragraphs[0]; p.text = title
    p.font.name = FONT_H; p.font.size = Pt(16); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6)
    p = tf.add_paragraph(); p.text = desc
    p.font.name = FONT_B; p.font.size = Pt(10); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER

def ns():
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        sp = ph._element; sp.getparent().remove(sp)
    return s

# ═══════════════════════════════════════
# SLIDE A — 项目背景与定位
# ═══════════════════════════════════════
s = ns()
# Title area
tb(s, LM, Emu(300000), Emu(8300000), Emu(700000), '需求分析', FONT_H, Pt(44), INK, True)
tb(s, LM, Emu(850000), Emu(8300000), Emu(400000), '项目背景与定位', FONT_B, Pt(18), MUTED)
# Divider line
rect(s, LM, Emu(1300000), Emu(1600000), Emu(5000), fill=ACC, radius=0)

# Left: project name + positioning
card(s, LM, Emu(1500000), Emu(5000000), Emu(1600000),
    'GeoguessrEX = GeoGuessr + Explorer + Exchange',
    ['GeoGuessr：地理猜点核心玩法',
     'Explorer：空间探索，发现校园',
     'Exchange：内容分享，互动交流',
     '定位：以图寻游戏为载体，融合空间认知、',
     '校园导览与文化传播的轻量级应用'],
    title_color=ACC)
# Right: why campus
card(s, LM+Emu(5200000), Emu(1500000), Emu(3550000), Emu(1600000),
    '为什么聚焦武汉大学校园？',
    ['校园空间范围明确，建筑辨识度高',
     '大量地标蕴含丰富的历史文化故事',
     '新生、游客、在校生均有真实使用场景',
     '不做通用型 GeoGuessr，做有在地特色',
     '的"武大街景 Guessr"'],
    title_color=GRN)

# Bottom key insight
rect(s, LM, Emu(3350000), CW, Emu(600000), fill=ACC)
tb(s, LM+Emu(200000), Emu(3420000), CW-Emu(400000), Emu(180000),
    '核心思路', FONT_H, Pt(16), WHITE, True)
tb(s, LM+Emu(200000), Emu(3580000), CW-Emu(400000), Emu(300000),
    '不是单纯复刻一个猜地点小游戏，而是把空间认知、校园导览和文化传播结合起来，',
    FONT_B, Pt(14), WHITE)

# Target form
tb(s, LM, Emu(4100000), Emu(8300000), Emu(350000),
    '最终目标：一个聚焦武大场景、具备实际使用价值的图寻原型系统',
    FONT_B, Pt(15), MUTED)

# ═══════════════════════════════════════
# SLIDE B — 目标用户分析
# ═══════════════════════════════════════
s = ns()
tb(s, LM, Emu(300000), Emu(8300000), Emu(700000), '需求分析', FONT_H, Pt(44), INK, True)
tb(s, LM, Emu(850000), Emu(8300000), Emu(400000), '目标用户分析', FONT_B, Pt(18), MUTED)
rect(s, LM, Emu(1300000), Emu(1600000), Emu(5000), fill=ACC, radius=0)

users = [
    ('武汉大学新生', ACC,
     ['核心需求：尽快熟悉校园环境',
      '痛点：地图能解决"怎么走"，但不能解决"这是哪""长什么样""为什么重要"',
      '价值：游戏化方式帮助新生主动建立校园空间认知']),
    ('游客与访客', GRN,
     ['核心需求：游览体验 + 文化体验',
      '痛点：传统导览方式被动，互动性和趣味性有限',
      '价值：轻量互动中了解校园地标和建筑故事']),
    ('在校生与校友', RGBColor(0x7F,0x2D,0x15),
     ['核心需求：趣味体验 + 情感连接',
      '痛点：对熟悉校园缺乏重新发现的契机',
      '价值：通过游戏唤起校园记忆，增强文化认同']),
]
uw = Emu(2700000); uh = Emu(2800000); ugap = Emu(100000)
for i,(title,color,lines) in enumerate(users):
    ux = LM + i*(uw+ugap)
    # Header
    hb = rect(s, ux, Emu(1500000), uw, Emu(500000), fill=color)
    htf = hb.text_frame; htf.margin_top = Emu(30000)
    p = htf.paragraphs[0]; p.text = title
    p.font.name = FONT_H; p.font.size = Pt(18); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Content card
    card(s, ux, Emu(2000000), uw, Emu(2300000), '', lines, title_color=INK, title_size=Pt(12))

# Bottom summary
tb(s, LM, Emu(4500000), CW, Emu(350000),
    '三类用户，一个交集：都需要更主动、更有趣、更有内容的方式来认识和理解校园空间。',
    FONT_B, Pt(15), INK, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE C — 核心需求
# ═══════════════════════════════════════
s = ns()
tb(s, LM, Emu(300000), Emu(8300000), Emu(700000), '需求分析', FONT_H, Pt(44), INK, True)
tb(s, LM, Emu(850000), Emu(8300000), Emu(400000), '核心功能需求', FONT_B, Pt(18), MUTED)
rect(s, LM, Emu(1300000), Emu(1600000), Emu(5000), fill=ACC, radius=0)

# Left column - basic needs
card(s, LM, Emu(1500000), Emu(3700000), Emu(3300000),
    '基础功能需求', [
    '① 题目展示：校园地点街景图或静态图片',
    '② 地图交互：在校园地图上点击标记猜测位置',
    '③ 结果判定：Haversine 距离计算 + 指数衰减评分',
    '④ 连续答题：多题串联，形成完整闯关体验',
    '⑤ 结算面板：总分、平均距离、历史对比'],
    title_color=ACC)

# Right column - unique feature
card(s, LM+Emu(3900000), Emu(1500000), Emu(4850000), Emu(1600000),
    '⭐ 核心特色：地点背景与故事展示', [
    '猜测完成后，不只告诉你"猜对没有""差了多少米"',
    '还会展示该地点的建筑介绍、历史背景、功能定位',
    '乃至校园小故事'],
    title_color=GRN)

# Bottom: feature highlight
card(s, LM+Emu(3900000), Emu(3350000), Emu(4850000), Emu(1400000),
    '', ['这个功能让系统从"图寻游戏"',
    '升级为：校园认知工具 + 文化传播载体'],
    title_color=ACC, title_size=Pt(14))

# Bottom note
tb(s, LM, Emu(4980000), CW, Emu(200000),
    '最终的"武大街景 Guessr"：以街景为线索、以地图猜点为核心玩法、以地点故事为内容补充。',
    FONT_B, Pt(12), MUTED)

# ═══════════════════════════════════════
# SLIDE D — 项目价值与意义
# ═══════════════════════════════════════
s = ns()
tb(s, LM, Emu(300000), Emu(8300000), Emu(700000), '需求分析', FONT_H, Pt(44), INK, True)
tb(s, LM, Emu(850000), Emu(8300000), Emu(400000), '项目价值与意义', FONT_B, Pt(18), MUTED)
rect(s, LM, Emu(1300000), Emu(1600000), Emu(5000), fill=ACC, radius=0)

values = [
    ('对新生的价值', ACC,
     '以更有趣、更有参与感的方式熟悉校园\n降低认路成本，缓解入学陌生感\n从"被动看地图"变为"主动探索空间"'),
    ('对游客的价值', GRN,
     '将传统静态导览转化为轻量化互动体验\n在"猜地点"和"看故事"中主动参与\n获取比传统导览更生动的校园文化信息'),
    ('对文化传播的价值', RGBColor(0x7F,0x2D,0x15),
     '校园建筑不仅是空间单元，更是历史文化载体\n连接地图、地点、建筑与故事\n让空间信息成为传播校园文化的媒介'),
]
vw = Emu(2750000); vh = Emu(2800000); vgap = Emu(80000)
for i,(title,color,desc) in enumerate(values):
    vx = LM + i*(vw+vgap)
    hb = rect(s, vx, Emu(1500000), vw, Emu(500000), fill=color)
    htf = hb.text_frame; htf.margin_top = Emu(30000)
    p = htf.paragraphs[0]; p.text = title
    p.font.name = FONT_H; p.font.size = Pt(18); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    card(s, vx, Emu(2000000), vw, Emu(2300000), '', desc.split('\n'), title_color=INK, title_size=Pt(12))

# Bottom definition
rect(s, LM, Emu(4500000), CW, Emu(500000), fill=ACC)
tb(s, LM+Emu(200000), Emu(4550000), CW-Emu(400000), Emu(200000),
    'GeoguessrEX = 游戏趣味性 + 校园认知工具 + 轻量化导览 + 校园文化展示', FONT_H, Pt(18), WHITE, True, PP_ALIGN.CENTER)

# ── Reorder: move the 4 new slides to the front ──
total = len(prs.slides)
new_indices = list(range(total-4, total))  # last 4 are new
sldIdLst = prs.slides._sldIdLst
entries = list(sldIdLst)
new_entries = [entries[i] for i in new_indices]
# Remove new entries from end
for i in reversed(new_indices):
    sldIdLst.remove(entries[i])
# Insert at beginning
for entry in reversed(new_entries):
    sldIdLst.insert(0, entry)

out = 'v2.0_with_req.pptx'
prs.save(out)
print(f'Done! {len(prs.slides)} slides (4 new + 9 original) → {out}')
print('Close v2.0.pptx in PowerPoint then rename this over it.')
