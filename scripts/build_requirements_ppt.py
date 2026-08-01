"""
Rebuild 需求分析.pptx from PRD.md — FIXED: add-then-delete to avoid zip name conflicts.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Constants ──
SLIDE_W = 9144000
SLIDE_H = 5143500
INK      = RGBColor(0x44, 0x37, 0x28)
MUTED    = RGBColor(0x83, 0x5E, 0x54)
CARD_BG  = RGBColor(0xEB, 0xE2, 0xE0)
CARD_BDR = RGBColor(0xD1, 0xC8, 0xC6)
ACCENT   = RGBColor(0xB4, 0x4D, 0x28)
GREEN    = RGBColor(0x24, 0x4C, 0x47)
WHITE    = RGBColor(0xFF, 0xF8, 0xEF)
FONT_TITLE = 'Crimson Pro Bold'
FONT_BODY  = 'Open Sans'
LM = Emu(400000)
CW = SLIDE_W - LM * 2

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
# Use blank layout (index 6 in default template) if available, else default
layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

# ── Helpers ──

def rect(slide, l, t, w, h, fill=None, border=CARD_BDR, bw=Emu(4763), radius=0.08):
    s = slide.shapes.add_shape(5, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.fill.solid(); s.line.width = bw
    else: s.line.fill.background()
    s.adjustments[0] = radius
    return s

def textbox(slide, l, t, w, h, text, font=FONT_BODY, size=Pt(12), color=INK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.name = font; p.font.size = size; p.font.color.rgb = color; p.font.bold = bold
    p.alignment = align
    return tb, tf

def add_badge(slide, l, t, w, h, text, fill=ACCENT):
    s = rect(slide, l, t, w, h, fill=fill, border=None, radius=0.35)
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = text
    p.font.name = FONT_BODY; p.font.size = Pt(11); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return s

def header(slide, badge, title, sub):
    add_badge(slide, Emu(496119), Emu(350000), Emu(500000), Emu(190000), badge)
    textbox(slide, Emu(1060000), Emu(374000), Emu(7500000), Emu(150000), title, FONT_TITLE, Pt(36), INK, True)
    if sub:
        textbox(slide, Emu(496119), Emu(560000), Emu(8151763), Emu(200000), sub, FONT_BODY, Pt(13), MUTED)

def card(slide, l, t, w, h, title, lines, title_color=INK, fill=CARD_BG):
    s = rect(slide, l, t, w, h, fill=fill)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(120000); tf.margin_right = Emu(120000)
    tf.margin_top = Emu(80000); tf.margin_bottom = Emu(80000)
    p = tf.paragraphs[0]; p.text = title
    p.font.name = FONT_TITLE; p.font.size = Pt(14); p.font.color.rgb = title_color; p.font.bold = True
    p.space_after = Pt(6)
    for line in lines:
        p = tf.add_paragraph(); p.text = line
        p.font.name = FONT_BODY; p.font.size = Pt(10); p.font.color.rgb = INK
        p.space_after = Pt(3)
    return s

def table_card(slide, l, t, w, h, headers, rows):
    s = rect(slide, l, t, w, h, fill=CARD_BG)
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(80000); tf.margin_right = Emu(80000)
    tf.margin_top = Emu(60000); tf.margin_bottom = Emu(60000)
    p = tf.paragraphs[0]
    p.text = '  │  '.join(headers)
    p.font.name = FONT_BODY; p.font.size = Pt(9); p.font.color.rgb = ACCENT; p.font.bold = True
    p.space_after = Pt(6)
    p = tf.add_paragraph(); p.text = '─' * 60
    p.font.name = FONT_BODY; p.font.size = Pt(6); p.font.color.rgb = CARD_BDR
    p.space_after = Pt(4)
    for row in rows:
        p = tf.add_paragraph()
        p.text = '  │  '.join(str(c) for c in row)
        p.font.name = FONT_BODY; p.font.size = Pt(9); p.font.color.rgb = INK
        p.space_after = Pt(3)
    return s

def tag(slide, l, t, text, fill=ACCENT):
    s = rect(slide, l, t, Emu(1000000), Emu(120000), fill=fill, border=None, radius=0.4)
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = text; p.font.name = FONT_BODY; p.font.size = Pt(7)
    p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

def new_slide():
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        sp = ph._element; sp.getparent().remove(sp)
    return s

# ═══════════════════════════════════════
# BUILD 14 NEW SLIDES (they'll be appended after the 15 old ones)
# ═══════════════════════════════════════

# ── SLIDE 1 — 封面 ──
s = new_slide()
rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=CARD_BG, border=None)
rect(s, Emu(0), Emu(0), SLIDE_W, Emu(2200000), fill=ACCENT, border=None)
textbox(s, LM, Emu(800000), Emu(8000000), Emu(600000), 'GeoguessrEX', FONT_TITLE, Pt(72), WHITE, True)
textbox(s, LM, Emu(1400000), Emu(8000000), Emu(400000), '需求分析文档', FONT_TITLE, Pt(36), WHITE)
textbox(s, LM, Emu(2500000), Emu(8000000), Emu(300000), '以图寻游戏为载体的社交软件平台', FONT_BODY, Pt(18), RGBColor(0xFF, 0xE0, 0xD0))
textbox(s, LM, Emu(3200000), Emu(4000000), Emu(200000), '版本 v0.2.0  ·  2026-07-08  ·  草稿', FONT_BODY, Pt(12), MUTED)
textbox(s, LM, Emu(3500000), Emu(6000000), Emu(200000), 'Vite + React 18  ·  Express  ·  高德地图  ·  Google 街景  ·  Haversine 算法', FONT_BODY, Pt(10), MUTED)

# ── SLIDE 2 — 项目背景 ──
s = new_slide()
header(s, '项目背景', '产品定位与愿景', 'PRD §1.1  |  从 MVP 工具到社交平台')
card(s, LM, Emu(950000), Emu(5100000), Emu(1200000), '产品定位',
    ['GeoguessrEX 是一款以图寻游戏为载体的社交软件',
     '在现有 MVP 单人玩法基础上，加入用户系统、社交功能、', '题库共建与管理后台，成为具备 UGC 生态和社交连接的平台'])
card(s, LM + Emu(5300000), Emu(950000), Emu(3400000), Emu(1200000), '技术栈',
    ['前端：Vite + React 18', '地图：高德 AMap JSAPI + Google Maps Embed API',
     '后端：Express (Node.js) RESTful API', '数据：本地 JSON → 迁移至数据库', '包管理：pnpm workspace (monorepo)'])
rect(s, LM, Emu(2400000), CW, Emu(800000), fill=ACCENT, border=None)
textbox(s, LM + Emu(200000), Emu(2480000), CW - Emu(400000), Emu(300000),
        '核心转变：从"工具型 MVP" → "社交平台"', FONT_TITLE, Pt(18), WHITE, True)
textbox(s, LM + Emu(200000), Emu(2780000), CW - Emu(400000), Emu(300000),
        'v0.1.0 只解决"能不能玩"  ·  v0.2+ 解决"谁在玩、怎么一起玩、谁在出题"', FONT_BODY, Pt(12), RGBColor(0xFF, 0xE0, 0xD0))
table_card(s, LM, Emu(3450000), CW, Emu(1500000),
    ['模块', 'v0.1.0 (现有)', 'v0.2+ (本次)'],
    [['用户', '本地 JSON + Bearer Token', '邮箱注册/登录 · JWT · 游客 · 个人主页'],
     ['题库', 'questions.json 手动编辑', 'UGC 上传 · URL 解析 · 图片录入 · 审核'],
     ['社交', '无', '好友 · 动态 · 五种排行榜'],
     ['管理', '无', '独立管理后台 (apps/admin)']])

# ── SLIDE 3 — 现有基础 ──
s = new_slide()
header(s, '现有基础', 'v0.1.0 已实现能力', 'PRD §1.2  |  MVP 阶段已完成的技术积累')
modules = [
    ('单人图寻', '题库→街景→选点\n→距离→分数', ACCENT),
    ('街景展示', 'Google Maps Embed API\n360° Street View iframe', GREEN),
    ('猜点地图', '高德 AMap JSAPI\n点击放置标记 + 坐标转换', ACCENT),
    ('结果地图', '双标记展示猜测点\n与答案点 + 连线对比', GREEN),
    ('坐标转换', 'WGS-84 ↔ GCJ-02\n中国境外自动检测', ACCENT),
    ('评分算法', 'Haversine + 指数衰减\n满分 5000 · decay 2000km', GREEN),
    ('题库', '本地 JSON 分组存储\n5 道示例题目', ACCENT),
]
cw = Emu(2400000); ch = Emu(800000); gap = Emu(150000)
for i, (title, desc, color) in enumerate(modules):
    col = i % 4; row = i // 4
    cx = LM + col * (cw + gap)
    cy = Emu(950000) + row * (ch + gap + Emu(200000))
    b = rect(s, cx, cy, cw, ch, fill=color, border=None)
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(60000); tf.margin_right = Emu(60000); tf.margin_top = Emu(30000)
    p = tf.paragraphs[0]; p.text = title
    p.font.name = FONT_BODY; p.font.size = Pt(11); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
    p = tf.add_paragraph(); p.text = desc
    p.font.name = FONT_BODY; p.font.size = Pt(8); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
textbox(s, LM, Emu(3700000), CW, Emu(350000),
        '技术栈：Vite + React 18 前端  ·  Express 后端  ·  高德 AMap JSAPI  ·  Google Maps Embed API  ·  本地 JSON 存储',
        FONT_BODY, Pt(10), MUTED, align=PP_ALIGN.CENTER)

# ── SLIDE 4 — 用户角色 ──
s = new_slide()
header(s, '用户角色', '四种角色与权限矩阵', 'PRD §2  |  从游客到管理员的完整权限体系')
roles = [
    ('游客', '每日 ≤3 局单人练习\n浏览社区内容\n无需注册，数据不持久', ACCENT),
    ('注册用户', '全部功能：\n挑战 · 上传题目 · 加好友\n发动态 · 排行榜', GREEN),
    ('审核员', '题库审核权限\n+ 注册用户全部权限\n可由管理员任命', RGBColor(0x7F, 0x2D, 0x15)),
    ('管理员', '管理后台全部权限：\n题库管理 · 用户管理\n系统配置 · 数据统计', RGBColor(0x3B, 0x3B, 0x3B)),
]
rw = Emu(1950000); rh = Emu(2400000); rgap = Emu(100000)
for i, (role, desc, color) in enumerate(roles):
    rx = LM + i * (rw + rgap)
    b = rect(s, rx, Emu(950000), rw, rh, fill=color, border=None)
    tf = b.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(60000); tf.margin_right = Emu(60000); tf.margin_top = Emu(60000)
    p = tf.paragraphs[0]; p.text = role
    p.font.name = FONT_TITLE; p.font.size = Pt(18); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(16)
    for line in desc.split('\n'):
        p = tf.add_paragraph(); p.text = line
        p.font.name = FONT_BODY; p.font.size = Pt(10); p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
textbox(s, LM, Emu(3550000), CW, Emu(400000),
        '权限递进：游客（少量体验）→ 注册用户（核心功能）→ 审核员（内容把关）→ 管理员（系统治理）',
        FONT_BODY, Pt(11), MUTED)
textbox(s, LM, Emu(3900000), CW, Emu(400000),
        '游客模式为新增概念（v0.1.0 无），降低首次使用门槛，注册后数据持久化并解锁全部功能。',
        FONT_BODY, Pt(10), MUTED)

# ── SLIDE 5 — 功能全景 ──
s = new_slide()
header(s, '功能需求', '五大功能模块全景', 'PRD §3  |  v0.2 迭代范围')
panels = [
    ('用户系统', ['邮箱注册/登录', '个人主页 + 数据面板', '游客模式（≤3局/天）', '忘记密码 / 记住我']),
    ('图寻挑战', ['单人练习（3/5/10题）', '限时挑战（60s倒计时）', '街景模式 + 图片模式', '结算面板 + 历史对比']),
    ('题库系统', ['URL 智能解析导入', '图片手动录入', '四步审核流程', '我的题库管理']),
    ('社交功能', ['好友系统（搜索/申请/同意）', '社区动态（分享+点赞+评论）', '五种排行榜', '个人主页']),
    ('管理后台', ['题库审核队列（批量操作）', '用户管理（封禁/角色分配）', '数据统计（平台+用户趋势）', '系统配置（可调参数面板）']),
]
pw = Emu(1680000); ph_ = Emu(2100000); pgap = Emu(50000)
for i, (title, lines) in enumerate(panels):
    px = LM + i * (pw + pgap)
    rect(s, px, Emu(950000), pw, Emu(350000), fill=ACCENT if i % 2 == 0 else GREEN, border=None)
    textbox(s, px + Emu(40000), Emu(960000), pw - Emu(80000), Emu(330000),
            title, FONT_TITLE, Pt(14), WHITE, True, PP_ALIGN.CENTER)
    cb = rect(s, px, Emu(1300000), pw, Emu(1750000), fill=CARD_BG)
    ctf = cb.text_frame; ctf.word_wrap = True
    ctf.margin_left = Emu(60000); ctf.margin_right = Emu(40000); ctf.margin_top = Emu(50000)
    for j, line in enumerate(lines):
        p = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
        p.text = f'• {line}'; p.font.name = FONT_BODY; p.font.size = Pt(10)
        p.font.color.rgb = INK; p.space_after = Pt(8)

# ── SLIDE 6 — 用户系统 ──
s = new_slide()
header(s, '功能需求', '用户系统', 'PRD §3.1  |  注册 / 登录 / 个人主页 / 游客模式')
card(s, LM, Emu(950000), Emu(4150000), Emu(1400000), '注册',
    ['邮箱 + 密码注册', '必填：邮箱、密码、昵称', '密码 ≥8 位，含字母和数字',
     '选填头像', '邮箱验证码（6位数字）', '注册后自动登录'], title_color=ACCENT)
card(s, LM + Emu(4350000), Emu(950000), Emu(4400000), Emu(1400000), '登录',
    ['邮箱 + 密码登录', '"记住我"：7天免登录（refresh token）',
     '密码错误 ≥5 次 → 锁定 15 分钟', '忘记密码 → 邮箱重置链接'], title_color=GREEN)
card(s, LM, Emu(2600000), Emu(4150000), Emu(1600000), '个人主页',
    ['头像 · 昵称 · 注册时间 · 简介', '数据面板：总局数、平均分、最佳成绩、上传数',
     '动态列表（时间倒序）', '编辑个人信息入口'], title_color=ACCENT)
card(s, LM + Emu(4350000), Emu(2600000), Emu(4400000), Emu(1600000), '游客模式',
    ['每日 ≤3 局单人练习', '可浏览社区内容', '不可加好友、发动态、上传题目',
     '成绩不纳入排行榜', '数据不持久'], title_color=GREEN)
textbox(s, LM, Emu(4450000), CW, Emu(300000),
        '安全策略：bcrypt 加盐哈希  ·  JWT access token (2h) + refresh token (7天)  ·  注册/登录接口限流',
        FONT_BODY, Pt(10), MUTED)

# ── SLIDE 7 — 图寻挑战 ──
s = new_slide()
header(s, '功能需求', '图寻挑战', 'PRD §3.2  |  单人练习 / 限时挑战 / 题目展示方式')
col_w = Emu(4200000)
card(s, LM, Emu(950000), col_w, Emu(2000000), '单人练习',
    ['从已审核公共题库随机抽题', '无时间限制，自由观察和判断', '每局可选 3 / 5 / 10 题',
     '提交后展示：双标记地图 + 距离 + 单题得分', '结算面板：总分、平均距离、最佳单题、历史对比',
     '成绩写入个人统计，可选分享至社区'], title_color=ACCENT)
card(s, LM + col_w + Emu(150000), Emu(950000), Emu(4400000), Emu(2000000), '限时挑战',
    ['从已审核公共题库随机抽题', '每道题 60s 倒计时（管理后台可配）',
     '超时 → 该题 0 分，自动下一题', '每局固定 5 题', '每日 3 次免费机会',
     '独立排行榜（与练习模式分开）'], title_color=GREEN)
card(s, LM, Emu(3200000), CW, Emu(1600000), '题目展示方式：街景模式 + 图片模式', [], title_color=INK)
table_card(s, LM, Emu(3550000), CW, Emu(1100000),
    ['模式', '数据来源', '展示形式', '适用场景'],
    [['街景模式', 'Google Maps 街景', 'iframe 嵌入 360° 全景', 'Google 街景覆盖区域'],
     ['图片模式', '用户上传的静态图片', '图片全屏 / 大图展示', '无街景覆盖地区、特色场景']])

# ── SLIDE 8 — 题库系统（上）数据模型 ──
s = new_slide()
header(s, '功能需求', '题库系统 — 数据模型', 'PRD §3.3.1  |  街景题 / 图片题 / 审核信息')
card(s, LM, Emu(950000), Emu(2750000), Emu(2400000), '基础信息（所有题目）',
    ['id: uuid', 'type: "streetview" | "image"', 'title: string（必填）',
     'description: string（出题者内部备注）', 'tags: string[]（国家/城市/场景类型）',
     'difficulty: easy/medium/hard（系统自动计算）', 'createdAt: timestamp'], title_color=INK)
card(s, LM + Emu(2900000), Emu(950000), Emu(2900000), Emu(2400000), '街景坐标 (type=streetview)',
    ['lat/lng: number（WGS-84 必填）', 'heading: number（0°~360°）',
     'pitch: number（-90°~90°）', 'fov: number（默认 100）', 'panoId: string（Google 全景 ID）'],
    title_color=ACCENT)
card(s, LM + Emu(6000000), Emu(950000), Emu(2750000), Emu(2400000), '图片坐标 (type=image)',
    ['imageUrl: string（存储 URL）', 'imageLat / imageLng: number', '图片格式：JPG/PNG/WebP ≤10MB',
     '建议 ≥1920×1080', '同一地点可上传多张不同角度'], title_color=GREEN)
table_card(s, LM, Emu(3600000), CW, Emu(1200000),
    ['审核字段', '类型', '说明'],
    [['authorId', 'uuid', '上传者用户 ID'],
     ['status', 'pending/approved/rejected', '审核状态（自动初始化）'],
     ['reviewerId', 'uuid', '审核员用户 ID'],
     ['reviewNote', 'string', '审核备注（驳回时必填）']])

# ── SLIDE 9 — 题库系统（下）上传与审核 ──
s = new_slide()
header(s, '功能需求', '题库系统 — 上传与审核', 'PRD §3.3.2-5  |  URL 导入 / 手动录入 / 审核流程 / 我的题库')
card(s, LM, Emu(950000), Emu(4200000), Emu(1900000), 'URL 导入',
    ['粘贴 Google Maps 街景链接', '自动解析经纬度、heading、', 'pitch、fov、panoId 等字段',
     '支持 3 种 URL 格式', '用户只需补充标题和标签'], title_color=ACCENT)
card(s, LM + Emu(2100000), Emu(950000), Emu(1600000), Emu(1900000), '手动录入',
    ['手动输入 WGS-84 经纬度', '上传图片（JPG/PNG/WebP）', '填写标题 + 选填标签', '适用于非街景场景'],
    title_color=GREEN)
card(s, LM + Emu(3900000), Emu(950000), Emu(2200000), Emu(1900000), '审核流程（5步）',
    ['① 提交 → pending', '② 审核员认领', '③ 查验（坐标/图片/合规）',
     '④ 判定：通过 → approved', '   驳回 → rejected + 原因', '⑤ 驳回后用户可修改重提'], title_color=INK)
card(s, LM + Emu(6300000), Emu(950000), Emu(2450000), Emu(1900000), '我的题库',
    ['查看自己上传的全部题目', '按状态筛选（全部/待审/已通过/已驳回）',
     '已驳回题目可编辑后重提', '已通过题目可删除（从公共库移除）'], title_color=GREEN)
for i, (label, color) in enumerate([('pending\n待审核', ACCENT), ('approved\n已通过', GREEN), ('rejected\n已驳回', RGBColor(0x8C, 0x3B, 0x2A))]):
    s2 = rect(s, LM + i * Emu(2800000), Emu(3100000), Emu(2400000), Emu(600000), fill=color, border=None, radius=0.1)
    tf = s2.text_frame; tf.word_wrap = True; tf.margin_left = Emu(60000); tf.margin_right = Emu(60000); tf.margin_top = Emu(40000)
    p = tf.paragraphs[0]; p.text = label; p.font.name = FONT_BODY; p.font.size = Pt(10)
    p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
textbox(s, LM + Emu(2400000), Emu(3200000), Emu(200000), Emu(200000), '→', FONT_BODY, Pt(20), MUTED, align=PP_ALIGN.CENTER)
textbox(s, LM + Emu(5200000), Emu(3200000), Emu(200000), Emu(200000), '→', FONT_BODY, Pt(20), MUTED, align=PP_ALIGN.CENTER)
for i, pill in enumerate(['坐标匹配', '坐标精度', '图片质量', '内容合规']):
    tag(s, LM + i * Emu(1050000), Emu(3900000), pill, ACCENT)
textbox(s, LM, Emu(4050000), CW, Emu(200000), '审核标准：以上四项全部通过方可通过审核', FONT_BODY, Pt(9), MUTED)

# ── SLIDE 10 — 社交功能 ──
s = new_slide()
header(s, '功能需求', '社交功能', 'PRD §3.4  |  好友系统 / 社区动态 / 排行榜')
card(s, LM, Emu(950000), Emu(3000000), Emu(2300000), '好友系统',
    ['搜索用户（昵称/ID）→ 发送申请', '→ 对方通知 → 同意/拒绝 → 互为好友',
     '好友列表：头像·昵称·活跃时间·统计', '预留：好友对战（后续迭代）'], title_color=ACCENT)
card(s, LM + Emu(2500000), Emu(950000), Emu(3000000), Emu(2400000), '社区动态',
    ['完成一局后 →「分享战绩」', '动态内容：战绩卡片 + 文字感言', '全站动态流（按时间倒序）',
     '互动：点赞 + 评论'], title_color=GREEN)
card(s, LM + Emu(5700000), Emu(950000), Emu(3050000), Emu(2400000), '排行榜（5种）',
    ['总积分榜 · 不重置', '周榜 · 每周一 00:00 重置', '日榜 · 每日 00:00 重置',
     '限时挑战榜 · 独立不重置', '好友榜 · 仅展示好友排名'], title_color=INK)
textbox(s, LM, Emu(3550000), CW, Emu(300000),
        '排行榜展示：排名、头像、昵称、分数、总局数、平均距离  |  游客成绩不纳入任何排行榜',
        FONT_BODY, Pt(10), MUTED)
textbox(s, LM, Emu(3850000), CW, Emu(400000),
        '设计思路：游戏是连接器，社交是留存引擎。一个人的探索变成一群人的分享与竞技。',
        FONT_BODY, Pt(11), INK, True)

# ── SLIDE 11 — 管理后台 ──
s = new_slide()
header(s, '功能需求', 'Web 管理后台', 'PRD §3.5  |  独立前端应用 (apps/admin)，管理员 & 审核员可访问')
admin_modules = [
    ('题库管理', ['审核队列：待审列表 + 预览 + 一键通过/驳回', '批量操作：勾选多题，批量通过或驳回',
                 '题库列表：按状态/类型/上传者/标签筛选', '编辑 / 下架 / 删除题目'], ACCENT),
    ('用户管理', ['用户列表：搜索 + 按角色筛选 + 按时间排序', '用户详情：个人信息 + 统计数据 + 最近活动',
                 '封禁/解封：封禁后无法登录和使用任何功能', '角色设置：普通用户 / 审核员 / 管理员'], GREEN),
    ('数据统计', ['平台概览：总用户数、总题目数、今日局数', '题库统计：待审/通过/驳回数、审核通过率',
                 '用户趋势：日新增注册、日活跃', '题目排行：热门题、高分题'], ACCENT),
    ('系统配置', ['限时挑战倒计时（默认 60s）', '游客每日局数（默认 3 局）',
                 '限时挑战每日次数（默认 3 次）', '评分满分 & 衰减常数', '练习默认题数（默认 5 题）'], GREEN),
]
mw = Emu(4200000); mh = Emu(1800000); mgap = Emu(150000)
for i, (title, lines, color) in enumerate(admin_modules):
    col = i % 2; row = i // 2
    mx = LM + col * (mw + mgap); my = Emu(950000) + row * (mh + mgap)
    card(s, mx, my, mw, mh, title, lines, title_color=color)

# ── SLIDE 12 — 非功能需求 ──
s = new_slide()
header(s, '非功能需求', '性能 · 安全 · 可维护性 · 兼容性', 'PRD §4')
nfr_cards = [
    ('性能', ACCENT, ['API 响应时间 P95 < 500ms', '图片首屏加载 < 2s（CDN）', '并发在线用户 ≥ 500', '图片上传响应 < 3s']),
    ('安全', GREEN, ['bcrypt 加盐哈希密码', 'JWT access(2h) + refresh(7d)', '文件上传：扩展名+MIME+大小校验',
                    '防 XSS · CSRF · SQL 注入', '接口限流：注册/登录 IP 每分钟 ≤10 次']),
    ('可维护性', INK, ['RESTful API 规范', '/api/ 用户端 · /api/admin/ 管理端', '数据库 migration 脚本 + 版本回滚',
                      '核心逻辑单元测试', '日志：请求 + 错误 + 审核操作']),
    ('兼容性', ACCENT, ['Chrome / Firefox / Edge 最新两个大版本', '桌面端最低 1280×720']),
]
nw = Emu(4200000); nh = Emu(1700000); ngap = Emu(150000)
for i, (title, color, lines) in enumerate(nfr_cards):
    col = i % 2; row = i // 2
    nx = LM + col * (nw + ngap); ny = Emu(950000) + row * (nh + ngap)
    card(s, nx, ny, nw, nh, title, lines, title_color=color)

# ── SLIDE 13 — 技术建议 ──
s = new_slide()
header(s, '技术建议', '数据库 · 存储 · 目录 · 模块划分', 'PRD §5')
card(s, LM, Emu(950000), Emu(4200000), Emu(1800000), '数据库',
    ['开发：SQLite（零配置、快速迭代）', '生产：PostgreSQL + PostGIS', '（地理空间查询支持："附近题目""区域筛选"）'],
    title_color=ACCENT)
card(s, LM + Emu(2000000), Emu(950000), Emu(1800000), Emu(1800000), '文件存储',
    ['开发：本地文件系统', '（apps/api/uploads/）', '生产：对象存储（OSS/S3）', '+ CDN 加速'], title_color=GREEN)
card(s, LM + Emu(4000000), Emu(950000), Emu(1600000), Emu(1800000), '目录结构',
    ['apps/api/  后端 API（现有）', 'apps/web/  用户端前端（现有）', 'apps/admin/ 管理后台（新增）'], title_color=INK)
card(s, LM, Emu(2500000), CW, Emu(2450000), '后端模块划分（路由层 apps/api/src/routes/）', [], title_color=INK)
table_card(s, LM, Emu(2880000), CW, Emu(1980000),
    ['路由文件', '职责', '主要接口'],
    [['auth.js', '注册/登录/登出/密码重置', 'POST /api/auth/*'],
     ['users.js', '用户信息查询与修改', 'GET/PUT /api/users/:id'],
     ['questions.js', '题库查询', 'GET /api/questions'],
     ['upload.js', '题目上传（URL解析+图片）', 'POST /api/questions'],
     ['game.js', '回合提交与会话', 'POST /api/submit'],
     ['social.js', '好友/动态/评论/点赞', '好友CRUD、动态发布/流'],
     ['leaderboard.js', '排行榜查询', 'GET /api/leaderboard?type='],
     ['admin/*.js', '管理端操作', '审核/用户管理/统计/配置']])

# ── SLIDE 14 — 版本规划 ──
s = new_slide()
header(s, '版本规划', 'v0.2 → v1.0 路线图', 'PRD §6  |  五个阶段交付')
phases = [
    ('v0.2', '用户系统\n+ 题库升级', ACCENT, ['邮箱注册/登录', 'JWT 鉴权', 'Google Maps URL 导入', '图片手动录入', '审核流程', '数据库迁移']),
    ('v0.3', '社交核心', GREEN, ['好友系统', '个人主页', '社区动态（分享+点赞+评论）', '五种排行榜']),
    ('v0.4', '限时挑战\n+ 管理后台', RGBColor(0x7F, 0x2D, 0x15), ['限时挑战模式', '独立排行榜', '审核队列', '用户管理', '数据统计', '系统配置']),
    ('v0.5', '体验优化', RGBColor(0x83, 0x5E, 0x54), ['CDN + 缓存', '游客模式', '性能调优']),
    ('v1.0', '正式版', RGBColor(0x3B, 0x3B, 0x3B), ['全面测试与缺陷修复', '生产环境部署', '运营数据看板']),
]
pw_ = Emu(1650000); ph__ = Emu(2800000)
for i, (ver, title, color, items) in enumerate(phases):
    px = LM + i * (pw_ + Emu(50000))
    hb = rect(s, px, Emu(950000), pw_, Emu(550000), fill=color, border=None)
    htf = hb.text_frame; htf.word_wrap = True; htf.margin_top = Emu(30000)
    p = htf.paragraphs[0]; p.text = ver
    p.font.name = FONT_TITLE; p.font.size = Pt(24); p.font.color.rgb = WHITE; p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p = htf.add_paragraph(); p.text = title
    p.font.name = FONT_BODY; p.font.size = Pt(9); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    cb = rect(s, px, Emu(1500000), pw_, Emu(2250000), fill=CARD_BG)
    ctf = cb.text_frame; ctf.word_wrap = True
    ctf.margin_left = Emu(50000); ctf.margin_right = Emu(30000); ctf.margin_top = Emu(50000)
    for j, item in enumerate(items):
        p = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
        p.text = f'• {item}'; p.font.name = FONT_BODY; p.font.size = Pt(10)
        p.font.color.rgb = INK; p.space_after = Pt(7)
    if i < len(phases) - 1:
        textbox(s, px + pw_ + Emu(10000), Emu(1100000), Emu(30000), Emu(300000),
                '→', FONT_BODY, Pt(20), ACCENT, True, PP_ALIGN.CENTER)
textbox(s, LM, Emu(3950000), CW, Emu(300000),
        'v0.2 为当前迭代版本，v1.0 为正式发布版。每个版本约 2-3 周，总计约 10-15 周完成全部功能。',
        FONT_BODY, Pt(11), MUTED)
textbox(s, LM, Emu(4250000), CW, Emu(300000),
        'MVP 已验证核心玩法可行，v0.2 起补齐用户体系与 UGC 能力，v0.3 起构建社交网络效应。',
        FONT_BODY, Pt(11), INK, False, PP_ALIGN.CENTER)

# ── Save ──
prs.save('需求分析.pptx')
print(f'Done! {len(prs.slides)} clean slides written to 需求分析.pptx')
