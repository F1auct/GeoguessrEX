"""
Add 4 core code logic slides to v1.4.pptx after slide 7 (技术架构).
Slides: 数学公式, 坐标转换, 判分流水线, 游戏状态机
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy
from lxml import etree

# ── Constants ──
SLIDE_W = 9144000
SLIDE_H = 5143500

# Colors (matching existing PPT)
INK = RGBColor(0x44, 0x37, 0x28)
MUTED = RGBColor(0x83, 0x5E, 0x54)
CARD_BG = RGBColor(0xEB, 0xE2, 0xE0)
CARD_BORDER = RGBColor(0xD1, 0xC8, 0xC6)
ACCENT = RGBColor(0xB4, 0x4D, 0x28)
ACCENT_DARK = RGBColor(0x7F, 0x2D, 0x15)
GREEN = RGBColor(0x24, 0x4C, 0x47)

# Code theme
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
CODE_KEYWORD = RGBColor(0x56, 0x9C, 0xD6)  # blue
CODE_FUNC = RGBColor(0xDC, 0xDC, 0xAA)      # yellow
CODE_STRING = RGBColor(0xCE, 0x91, 0x78)     # orange
CODE_COMMENT = RGBColor(0x6A, 0x99, 0x5B)    # green
CODE_NUMBER = RGBColor(0xB5, 0xCE, 0xA8)     # light green
CODE_DEFAULT = RGBColor(0xD4, 0xD4, 0xD4)    # light gray

FONT_TITLE = 'Crimson Pro Bold'
FONT_BODY = 'Open Sans'
FONT_CODE = 'Consolas'
FONT_CODE_FALLBACK = 'Courier New'

prs = Presentation('v1.4.pptx')


# ── Helper functions ──

def add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    # Round corners
    shape.adjustments[0] = 0.08
    return shape


def add_textbox(slide, left, top, width, height, text="", font_name=FONT_BODY,
                font_size=Pt(14), color=INK, bold=False, alignment=PP_ALIGN.LEFT,
                word_wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox, tf


def add_code_block(slide, left, top, width, height):
    """Add a dark code panel, return the shape and text frame."""
    shape = add_rounded_rect(slide, left, top, width, height,
                             fill_color=CODE_BG, border_color=RGBColor(0x3C, 0x3C, 0x3C),
                             border_width=Emu(4763))
    shape.adjustments[0] = 0.04
    tf = shape.text_frame
    tf.word_wrap = True
    # Padding
    tf.margin_left = Emu(120000)
    tf.margin_right = Emu(120000)
    tf.margin_top = Emu(100000)
    tf.margin_bottom = Emu(100000)
    return shape, tf


def add_code_line(tf, text, color=CODE_DEFAULT, size=Pt(10), bold=False, indent=0):
    """Add a line of code text. First call uses first paragraph, subsequent calls add new paras."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.name = FONT_CODE
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(2)
    if indent:
        p.level = indent
    # Force monospace ascii font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), FONT_CODE)
    return p


def add_badge(slide, left, top, width, height, fill_color=ACCENT):
    """Add a colored badge pill (like existing '技术架构' badge)."""
    shape = add_rounded_rect(slide, left, top, width, height,
                             fill_color=fill_color)
    shape.adjustments[0] = 0.35
    return shape


def add_slide_header(slide, badge_text, title_text, subtitle_text):
    """Add consistent slide header: badge + title + subtitle."""
    # Badge
    badge_w = Emu(500000)
    badge_h = Emu(190000)
    badge = add_badge(slide, Emu(496119), Emu(389855), badge_w, badge_h, fill_color=ACCENT)
    btf = badge.text_frame
    btf.word_wrap = False
    bp = btf.paragraphs[0]
    bp.text = badge_text
    bp.font.name = FONT_BODY
    bp.font.size = Pt(11)
    bp.font.color.rgb = RGBColor(0xFF, 0xF8, 0xEF)
    bp.font.bold = True
    bp.alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide, Emu(1060000), Emu(414412), Emu(7000000), Emu(150000),
                text=title_text, font_name=FONT_TITLE, font_size=Pt(40),
                color=INK, bold=True)

    # Subtitle
    add_textbox(slide, Emu(496119), Emu(580000), Emu(8151763), Emu(250000),
                text=subtitle_text, font_name=FONT_BODY, font_size=Pt(14),
                color=MUTED, bold=False)


# ── Build slides ──

# We'll insert after slide 7 (index 6), so new slides go at indices 7, 8, 9, 10
INSERT_AFTER = 6  # 0-based index of slide 7 (技术架构)

def clone_slide_background(source_slide_idx, target_slide):
    """Copy background from an existing slide."""
    # Just use a blank slide - backgrounds are typically white in this template
    pass

# ────────────────────────────────────────────
# SLIDE A: 数学公式 — Haversine + 指数衰减评分
# ────────────────────────────────────────────
slide_a = prs.slides.add_slide(prs.slide_layouts[0])  # Use DEFAULT layout
# Remove any placeholder shapes
for ph in list(slide_a.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

add_slide_header(slide_a, "核心算法", "数学公式", "Haversine 大圆距离 + 指数衰减评分")

# ── Haversine panel (left) ──
left_margin = Emu(350000)
code_w = Emu(4100000)
code_h = Emu(3400000)
gap = Emu(200000)

hav_shape, hav_tf = add_code_block(slide_a, left_margin, Emu(950000), code_w, Emu(3700000))

add_code_line(hav_tf, "// Haversine 大圆距离公式", CODE_COMMENT, size=Pt(9))
add_code_line(hav_tf, "export function haversineDistanceKm(", CODE_KEYWORD, size=Pt(9))
add_code_line(hav_tf, "  lat1, lng1, lat2, lng2)", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "  const R = 6371; // 地球半径 (km)", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "  const dLat = toRad(lat2 - lat1);", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "  const dLng = toRad(lng2 - lng1);", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "  const a = sin²(dLat/2) +", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "    cos(lat1)·cos(lat2)·sin²(dLng/2);", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "  const c = 2·atan2(√a, √(1−a));", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "  return R * c;  // 球面最短距离", CODE_DEFAULT, size=Pt(9))
add_code_line(hav_tf, "}", CODE_DEFAULT, size=Pt(9))

add_code_line(hav_tf, "", CODE_DEFAULT, size=Pt(6))
add_code_line(hav_tf, "// 算法来源: utils/haversine.js", CODE_COMMENT, size=Pt(8))

# ── Annotation box for Haversine ──
anno_w = Emu(4200000)
add_textbox(slide_a, left_margin, Emu(4750000), anno_w, Emu(250000),
            text="球面三角学经典公式，计算两点间最短弧长。输入 WGS-84 坐标，输出千米距离。",
            font_name=FONT_BODY, font_size=Pt(9), color=MUTED)

# ── Scoring panel (right) ──
score_left = left_margin + code_w + gap + Emu(10000)
score_shape, score_tf = add_code_block(slide_a, score_left, Emu(950000), code_w, Emu(2600000))

add_code_line(score_tf, "// 指数衰减评分公式", CODE_COMMENT, size=Pt(9))
add_code_line(score_tf, "export function scoreFromDistance(", CODE_KEYWORD, size=Pt(9))
add_code_line(score_tf, "  distanceKm)", CODE_DEFAULT, size=Pt(9))
add_code_line(score_tf, "  const maxScore = 5000;", CODE_NUMBER, size=Pt(9))
add_code_line(score_tf, "  const decay  = 2000;", CODE_NUMBER, size=Pt(9))
add_code_line(score_tf, "", CODE_DEFAULT, size=Pt(9))
add_code_line(score_tf, "  return Math.max(0,", CODE_DEFAULT, size=Pt(9))
add_code_line(score_tf, "    Math.round(", CODE_DEFAULT, size=Pt(9))
add_code_line(score_tf, "      maxScore *", CODE_DEFAULT, size=Pt(9))
add_code_line(score_tf, "      Math.exp(-distanceKm / decay)", CODE_DEFAULT, size=Pt(9))
add_code_line(score_tf, "    )", CODE_DEFAULT, size=Pt(9))
add_code_line(score_tf, "  );", CODE_DEFAULT, size=Pt(9))
add_code_line(score_tf, "}", CODE_DEFAULT, size=Pt(9))

add_code_line(score_tf, "", CODE_DEFAULT, size=Pt(6))
add_code_line(score_tf, "// 算法来源: utils/scoring.js", CODE_COMMENT, size=Pt(8))

# ── Formula annotation ──
add_textbox(slide_a, score_left, Emu(3650000), anno_w, Emu(250000),
            text="Score = max(0, round(5000 × e^(−d/2000)))",
            font_name=FONT_CODE, font_size=Pt(10), color=ACCENT_DARK, bold=True)

add_textbox(slide_a, score_left, Emu(3900000), anno_w, Emu(300000),
            text="d = 0 → 5000 分 (满分)    d = 2000 km → ≈ 1840 分    d → ∞ → 0 分",
            font_name=FONT_BODY, font_size=Pt(9), color=MUTED)

# ── Curve illustration note ──
add_textbox(slide_a, score_left, Emu(4300000), anno_w, Emu(350000),
            text="核心思想：距离越近得分越高，指数衰减确保远端仍有区分度。满分 5000，猜对位置即满分。",
            font_name=FONT_BODY, font_size=Pt(9), color=MUTED)

# ────────────────────────────────────────────
# SLIDE B: 坐标转换 — WGS-84 ↔ GCJ-02
# ────────────────────────────────────────────
slide_b = prs.slides.add_slide(prs.slide_layouts[0])
for ph in list(slide_b.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

add_slide_header(slide_b, "核心算法", "坐标转换", "WGS-84 ↔ GCJ-02 双向转换 + 中国境外检测")

# ── Left: outOfChina + transform ──
left_code, left_tf = add_code_block(slide_b, left_margin, Emu(950000), Emu(3700000), Emu(3200000))

add_code_line(left_tf, "// 中国境外检测 — 不做转换", CODE_COMMENT, size=Pt(9))
add_code_line(left_tf, "function outOfChina(lng, lat) {", CODE_KEYWORD, size=Pt(9))
add_code_line(left_tf, "  return lng < 72.004 || lng > 137.8347", CODE_DEFAULT, size=Pt(9))
add_code_line(left_tf, "      || lat < 0.8293  || lat > 55.8271;", CODE_DEFAULT, size=Pt(9))
add_code_line(left_tf, "}", CODE_DEFAULT, size=Pt(9))
add_code_line(left_tf, "", CODE_DEFAULT, size=Pt(6))
add_code_line(left_tf, "// WGS-84 → GCJ-02 (高德地图使用)", CODE_COMMENT, size=Pt(9))
add_code_line(left_tf, "export function wgs84ToGcj02(lng, lat) {", CODE_KEYWORD, size=Pt(9))
add_code_line(left_tf, "  if (outOfChina(lng, lat))", CODE_DEFAULT, size=Pt(9))
add_code_line(left_tf, "    return { lng, lat };  // 境外不变", CODE_DEFAULT, size=Pt(9))
add_code_line(left_tf, "  // 非线性偏移计算 (国家测绘局算法)", CODE_COMMENT, size=Pt(9))
add_code_line(left_tf, "  let dLat = transformLat(lng-105, lat-35);", CODE_DEFAULT, size=Pt(9))
add_code_line(left_tf, "  let dLng = transformLng(lng-105, lat-35);", CODE_DEFAULT, size=Pt(9))
add_code_line(left_tf, "  // ... 椭球参数修正 ...", CODE_COMMENT, size=Pt(9))
add_code_line(left_tf, "  return { lng: lng+dLng, lat: lat+dLat };", CODE_DEFAULT, size=Pt(9))
add_code_line(left_tf, "}", CODE_DEFAULT, size=Pt(9))

# ── Right: gcj02ToWgs84 ──
right_code, right_tf = add_code_block(slide_b, left_margin + Emu(3700000) + gap, Emu(950000),
                                       Emu(4100000), Emu(2200000))

add_code_line(right_tf, "// GCJ-02 → WGS-84 (用户点击后还原)", CODE_COMMENT, size=Pt(9))
add_code_line(right_tf, "export function gcj02ToWgs84(lng, lat) {", CODE_KEYWORD, size=Pt(9))
add_code_line(right_tf, "  if (outOfChina(lng, lat))", CODE_DEFAULT, size=Pt(9))
add_code_line(right_tf, "    return { lng, lat };", CODE_DEFAULT, size=Pt(9))
add_code_line(right_tf, "", CODE_DEFAULT, size=Pt(6))
add_code_line(right_tf, "  // 逆向迭代: 利用正向函数反推", CODE_COMMENT, size=Pt(9))
add_code_line(right_tf, "  const converted = wgs84ToGcj02(lng, lat);", CODE_DEFAULT, size=Pt(9))
add_code_line(right_tf, "  return {", CODE_DEFAULT, size=Pt(9))
add_code_line(right_tf, "    lng: 2*lng - converted.lng,", CODE_DEFAULT, size=Pt(9))
add_code_line(right_tf, "    lat: 2*lat - converted.lat", CODE_DEFAULT, size=Pt(9))
add_code_line(right_tf, "  };", CODE_DEFAULT, size=Pt(9))
add_code_line(right_tf, "}", CODE_DEFAULT, size=Pt(9))

add_code_line(right_tf, "", CODE_DEFAULT, size=Pt(6))
add_code_line(right_tf, "// 算法来源:", CODE_COMMENT, size=Pt(8))
add_code_line(right_tf, "// services/coordTransform.js", CODE_COMMENT, size=Pt(8))

# ── Conversion chain diagram ──
chain_y = Emu(3300000)
add_textbox(slide_b, left_margin + Emu(3700000) + gap, chain_y, Emu(4100000), Emu(200000),
            text="完整转换链路", font_name=FONT_TITLE, font_size=Pt(14),
            color=INK, bold=True)

# Chain boxes
chain_items = [
    ("题库存储\n(WGS-84)", ACCENT),
    ("wgs84ToGcj02\n展示到高德地图", ACCENT_DARK),
    ("用户点击标记\n(GCJ-02)", GREEN),
    ("gcj02ToWgs84\n还原为 WGS-84", ACCENT_DARK),
    ("后端 Haversine\n计算 (WGS-84)", GREEN),
]

box_w = Emu(1350000)
box_h = Emu(580000)
chain_start_x = left_margin + Emu(3700000) + gap
box_gap = Emu(50000)
arrow_y = chain_y + Emu(480000)

total_chain_w = len(chain_items) * box_w + (len(chain_items) - 1) * box_gap
chain_start_x += (Emu(4100000) - total_chain_w) // 2

for i, (label, color) in enumerate(chain_items):
    bx = chain_start_x + i * (box_w + box_gap)
    by = arrow_y
    box = add_rounded_rect(slide_b, bx, by, box_w, box_h,
                           fill_color=color, border_color=None)
    box.adjustments[0] = 0.1
    btf = box.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.text = label
    bp.font.name = FONT_BODY
    bp.font.size = Pt(7)
    bp.font.color.rgb = RGBColor(0xFF, 0xF8, 0xEF)
    bp.font.bold = True
    bp.alignment = PP_ALIGN.CENTER

# Arrow text
add_textbox(slide_b, chain_start_x, arrow_y - Emu(180000), total_chain_w, Emu(150000),
            text="→  WGS-84  →  GCJ-02  →  GCJ-02  →  WGS-84  →  WGS-84",
            font_name=FONT_CODE, font_size=Pt(7), color=MUTED, alignment=PP_ALIGN.CENTER)

# Bottom note
add_textbox(slide_b, left_margin, Emu(4250000), Emu(7800000), Emu(350000),
            text="为什么需要：高德地图（AMap JSAPI）依法必须使用 GCJ-02 坐标系，而 Google 街景和后端距离计算使用 WGS-84。"
                 "国内坐标会偏移 100-700 米，必须转换；境外坐标 outOfChina() 直接透传不做转换。",
            font_name=FONT_BODY, font_size=Pt(9), color=MUTED)


# ────────────────────────────────────────────
# SLIDE C: 判分流水线 — gradeAnswer
# ────────────────────────────────────────────
slide_c = prs.slides.add_slide(prs.slide_layouts[0])
for ph in list(slide_c.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

add_slide_header(slide_c, "核心逻辑", "判分流水线", "一次 gradeAnswer() 调用串联全部后端逻辑")

# ── Pipeline visualization ──
# 5 pipeline stages
pipeline_y = Emu(1000000)
stage_w = Emu(1550000)
stage_h = Emu(700000)
stage_gap = Emu(80000)
pipe_start_x = Emu(350000)

stages = [
    ("① 查题", "getQuestionById()\n从 questions.json\n按 ID 查找题目", ACCENT),
    ("② 校验", "guess.lat / lng\n类型检查，NaN 过滤\n400 错误直接返回", ACCENT_DARK),
    ("③ 算距离", "haversineDistanceKm()\nguess vs answer\nWGS-84 球面距离", GREEN),
    ("④ 评分", "scoreFromDistance()\n指数衰减转换\n距离 → 0~5000 分", ACCENT_DARK),
    ("⑤ 组装结果", "返回完整结果包\nquestionId, score,\ndistanceKm, answer…", GREEN),
]

for i, (title, desc, color) in enumerate(stages):
    sx = pipe_start_x + i * (stage_w + stage_gap)
    box = add_rounded_rect(slide_c, sx, pipeline_y, stage_w, stage_h,
                           fill_color=color, border_color=None)
    box.adjustments[0] = 0.08
    btf = box.text_frame
    btf.word_wrap = True
    btf.margin_left = Emu(60000)
    btf.margin_right = Emu(60000)
    btf.margin_top = Emu(40000)

    bp = btf.paragraphs[0]
    bp.text = title
    bp.font.name = FONT_BODY
    bp.font.size = Pt(11)
    bp.font.color.rgb = RGBColor(0xFF, 0xF8, 0xEF)
    bp.font.bold = True
    bp.alignment = PP_ALIGN.CENTER

    bp2 = btf.add_paragraph()
    bp2.text = desc
    bp2.font.name = FONT_BODY
    bp2.font.size = Pt(7)
    bp2.font.color.rgb = RGBColor(0xFF, 0xF8, 0xEF)
    bp2.alignment = PP_ALIGN.CENTER
    bp2.space_before = Pt(4)

# Arrows between stages
for i in range(len(stages) - 1):
    ax = pipe_start_x + (i + 1) * stage_w + i * stage_gap + stage_gap // 2
    add_textbox(slide_c, ax - Emu(20000), pipeline_y + stage_h // 2 - Emu(60000),
                Emu(60000), Emu(120000), text="→",
                font_name=FONT_BODY, font_size=Pt(16), color=INK, bold=True,
                alignment=PP_ALIGN.CENTER)

# ── Code block below ──
code_block_y = Emu(2000000)
code_shape, code_tf = add_code_block(slide_c, Emu(350000), code_block_y,
                                      Emu(8450000), Emu(2800000))

add_code_line(code_tf, "// POST /api/submit  →  services/gameService.js", CODE_COMMENT, size=Pt(9))
add_code_line(code_tf, "export function gradeAnswer(questionId, guess) {", CODE_KEYWORD, size=Pt(9))
add_code_line(code_tf, "  const q = getQuestionById(questionId);    // ① 查题", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "  if (!q) return { error: '...', status: 404 };", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "", CODE_DEFAULT, size=Pt(5))
add_code_line(code_tf, "  if (typeof guess?.lat !== 'number'          // ② 校验", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "   || Number.isNaN(guess.lat))", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "    return { error: '...', status: 400 };", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "", CODE_DEFAULT, size=Pt(5))
add_code_line(code_tf, "  const km = haversineDistanceKm(            // ③ 算距离", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "    guess.lat, guess.lng,", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "    q.streetView.lat, q.streetView.lng);", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "", CODE_DEFAULT, size=Pt(5))
add_code_line(code_tf, "  return {                                    // ④⑤ 评分+组装", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "    questionId, title, description,", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "    guess, answer: { lat, lng },", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "    distanceKm: km,", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "    score: scoreFromDistance(km) };", CODE_FUNC, size=Pt(9))
add_code_line(code_tf, "}", CODE_DEFAULT, size=Pt(9))

# ── Bottom note ──
add_textbox(slide_c, Emu(350000), Emu(4880000), Emu(8450000), Emu(200000),
            text="一次 HTTP POST → 五步串联 → 前端拿到完整结果包（含距离、分数、正确答案），无需额外请求。",
            font_name=FONT_BODY, font_size=Pt(9), color=MUTED)


# ────────────────────────────────────────────
# SLIDE D: 游戏状态机 — GamePage 五态流转
# ────────────────────────────────────────────
slide_d = prs.slides.add_slide(prs.slide_layouts[0])
for ph in list(slide_d.placeholders):
    sp = ph._element
    sp.getparent().remove(sp)

add_slide_header(slide_d, "核心逻辑", "游戏状态机", "GamePage 五态流转 — 前端最复杂的交互调度中心")

# ── State diagram ──
# Layout: 5 states in a cycle, with transitions
state_radius = Emu(600000)
center_x = SLIDE_W // 2
center_y = Emu(2550000)

# Position 5 states in a pentagon
import math
states = [
    ("loading", "加载题库中...\n全屏 loading", ACCENT),
    ("ready", "街景 + 地图就绪\n等待用户标记位置", GREEN),
    ("submitting", "提交中...\n按钮禁用 + spinner", ACCENT_DARK),
    ("result", "结果地图 + 得分\n距离 + 地点介绍", ACCENT_DARK),
    ("next", "切换到下一题\n重置所有状态 → ready", GREEN),
]

state_positions = []
for i, (name, desc, color) in enumerate(states):
    angle = -math.pi / 2 + i * (2 * math.pi / 5)  # Start from top
    sx = center_x + int(Emu(2400000) * math.cos(angle))
    sy = center_y + int(Emu(1800000) * math.sin(angle))
    state_positions.append((sx, sy))

    # State circle
    circle = slide_d.shapes.add_shape(
        9,  # OVAL
        sx - state_radius // 2, sy - state_radius // 2,
        state_radius, state_radius
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Emu(40000)
    ctf.margin_right = Emu(40000)

    cp = ctf.paragraphs[0]
    cp.text = name
    cp.font.name = FONT_CODE
    cp.font.size = Pt(12)
    cp.font.color.rgb = RGBColor(0xFF, 0xF8, 0xEF)
    cp.font.bold = True
    cp.alignment = PP_ALIGN.CENTER

    cp2 = ctf.add_paragraph()
    cp2.text = desc
    cp2.font.name = FONT_BODY
    cp2.font.size = Pt(6)
    cp2.font.color.rgb = RGBColor(0xFF, 0xF8, 0xEF)
    cp2.alignment = PP_ALIGN.CENTER
    cp2.space_before = Pt(3)

# ── Transition arrows (using connector lines) ──
transitions = [
    (0, 1, "题库加载完成"),
    (1, 2, "用户点击提交"),
    (2, 3, "后端返回结果"),
    (3, 4, "点击下一题"),
    (4, 1, "状态重置"),
]

for from_idx, to_idx, label in transitions:
    fx, fy = state_positions[from_idx]
    tx, ty = state_positions[to_idx]

    # Draw arrow line
    connector = slide_d.shapes.add_connector(
        1,  # STRAIGHT connector
        fx, fy, tx, ty
    )
    connector.line.color.rgb = MUTED
    connector.line.width = Pt(1.5)
    connector.line.dash_style = 2  # dash

    # Label at midpoint
    mx = (fx + tx) // 2
    my = (fy + ty) // 2
    add_textbox(slide_d, mx - Emu(500000), my - Emu(140000), Emu(1000000), Emu(120000),
                text=label, font_name=FONT_BODY, font_size=Pt(6), color=MUTED,
                alignment=PP_ALIGN.CENTER)

# ── State code snippet (bottom) ──
code_y = Emu(3700000)
code_shape, code_tf = add_code_block(slide_d, Emu(350000), code_y,
                                      Emu(8450000), Emu(1250000))

add_code_line(code_tf, "// GamePage.jsx — 状态驱动渲染", CODE_COMMENT, size=Pt(9))
add_code_line(code_tf, "const [status, setStatus] = useState('loading');", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "// 状态流转:", CODE_COMMENT, size=Pt(9))
add_code_line(code_tf, "// loading → ready → submitting → result → ready (下一题)", CODE_COMMENT, size=Pt(9))
add_code_line(code_tf, "", CODE_DEFAULT, size=Pt(5))
add_code_line(code_tf, "if (status === 'loading')  return <Loading />;     // 加载题库", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "if (status === 'ready')    return <PlayUI />;       // 街景+猜点地图", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "if (status === 'submitting') return <PlayUI disabled />; // 提交中", CODE_DEFAULT, size=Pt(9))
add_code_line(code_tf, "if (status === 'result')   return <ResultUI />;      // 结果页", CODE_DEFAULT, size=Pt(9))

# ── Bottom note ──
add_textbox(slide_d, Emu(350000), Emu(5000000), Emu(8450000), Emu(200000),
            text="五态覆盖完整游戏生命周期：加载 → 交互 → 提交 → 反馈 → 循环。每态对应独立 UI 布局，状态切换驱动整个页面的重新渲染。",
            font_name=FONT_BODY, font_size=Pt(9), color=MUTED)


# ── Reorder slides: move new slides after slide 7 (index 6) ──
# New slides are currently at the end (indices -4, -3, -2, -1)
# We need to move them to indices 7, 8, 9, 10
total = len(prs.slides)
# The 4 new slides are at indices total-4, total-3, total-2, total-1
new_slide_indices = list(range(total - 4, total))

# Build XML element list and reorder
sldIdLst = prs.slides._sldIdLst
entries = list(sldIdLst)
# Remove the 4 new entries from the end
new_entries = [entries[i] for i in new_slide_indices]
# Remove them from original positions (going backwards to preserve indices)
for i in reversed(new_slide_indices):
    sldIdLst.remove(entries[i])
# Insert them at position 7 (after the 7th slide, which is index 6)
insert_pos = 7  # After slide 7 (0-indexed 6)
for entry in reversed(new_entries):
    sldIdLst.insert(insert_pos, entry)

# Save
output_path = 'v1.4.pptx'
prs.save(output_path)
print(f'Done! Saved to {output_path}')
print(f'Total slides: {len(prs.slides)}')
for i, slide in enumerate(prs.slides):
    shapes = slide.shapes
    titles = []
    for s in shapes:
        if s.has_text_frame:
            t = s.text_frame.paragraphs[0].text[:50] if s.text_frame.paragraphs[0].text else ''
            if t:
                titles.append(t)
                break
    print(f'  Slide {i+1}: {titles[0] if titles else "(no text)"}')
